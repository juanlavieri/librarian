"""Office + PDF + image readers.

Each depends on an optional third-party library. When the dependency is missing
the reader degrades gracefully: it returns an empty block list (and the
registry simply skips that format) instead of crashing the whole build. Install
the ``documents`` extra to enable them all.
"""

from __future__ import annotations

import io
from typing import List

from ..models import Block


class PdfReader:
    extensions = frozenset({".pdf"})

    def parse(self, name: str, data: bytes) -> List[Block]:
        blocks: List[Block] = []
        try:
            from pypdf import PdfReader as _PdfReader  # type: ignore
        except Exception:
            try:
                from PyPDF2 import PdfReader as _PdfReader  # type: ignore
            except Exception:
                return blocks
        try:
            reader = _PdfReader(io.BytesIO(data))
            for i, page in enumerate(reader.pages, start=1):
                text = (page.extract_text() or "").strip()
                if text:
                    blocks.append(Block(type="page", text=text, location=f"p.{i}"))
        except Exception:
            return blocks
        return blocks


class DocxReader:
    extensions = frozenset({".docx"})

    def parse(self, name: str, data: bytes) -> List[Block]:
        try:
            import docx  # type: ignore
        except Exception:
            return []
        try:
            document = docx.Document(io.BytesIO(data))
        except Exception:
            return []
        paragraphs = [p.text for p in document.paragraphs if p.text and p.text.strip()]
        text = "\n".join(paragraphs).strip()
        if not text:
            return []
        return [Block(type="text", text=text, location=name)]


class PptxReader:
    extensions = frozenset({".pptx"})

    def parse(self, name: str, data: bytes) -> List[Block]:
        try:
            from pptx import Presentation  # type: ignore
        except Exception:
            return []
        try:
            prs = Presentation(io.BytesIO(data))
        except Exception:
            return []
        blocks: List[Block] = []
        for i, slide in enumerate(prs.slides, start=1):
            parts = [shp.text for shp in slide.shapes if hasattr(shp, "text") and shp.text]
            text = "\n".join(parts).strip()
            if text:
                blocks.append(Block(type="slide", text=text, location=f"slide {i}"))
        return blocks


class XlsxReader:
    extensions = frozenset({".xlsx", ".xlsm"})

    def parse(self, name: str, data: bytes) -> List[Block]:
        try:
            from openpyxl import load_workbook  # type: ignore
        except Exception:
            return []
        try:
            wb = load_workbook(io.BytesIO(data), read_only=True, data_only=True)
        except Exception:
            return []
        blocks: List[Block] = []
        for sheet in wb.worksheets:
            lines: List[str] = []
            for row in sheet.iter_rows(max_row=1000, values_only=True):
                cells = [str(c) for c in row if c is not None]
                if cells:
                    lines.append(" | ".join(cells))
            text = "\n".join(lines).strip()
            if text:
                blocks.append(Block(type="sheet", text=text, location=sheet.title))
        return blocks


class ImageOcrReader:
    extensions = frozenset({".png", ".jpg", ".jpeg", ".tiff", ".bmp", ".gif"})

    def parse(self, name: str, data: bytes) -> List[Block]:
        try:
            import pytesseract  # type: ignore
            from PIL import Image  # type: ignore
        except Exception:
            return []
        try:
            text = pytesseract.image_to_string(Image.open(io.BytesIO(data))).strip()
        except Exception:
            return []
        if not text:
            return []
        return [Block(type="text", text=text, location=name)]
